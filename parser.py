"""PPT text extraction using python-pptx library with layout analysis."""

from pathlib import Path

from pptx import Presentation


def extract_ppt(file_path: Path) -> list[dict]:
    """
    Extract text and layout metadata from a PowerPoint presentation.

    Args:
        file_path: Path to .pptx file

    Returns:
        List of slide dictionaries with:
        - slide_id: Unique slide identifier
        - slide_index: Zero-based slide number
        - paragraphs: List of paragraph objects with text, font size, position, type
        - layout_metadata: Summary of layout for each text element
        - slide_type: Classified slide type (title, body, agenda, closing, etc.)
    """
    presentation = Presentation(str(file_path))
    extracted = []

    for slide_index, slide in enumerate(presentation.slides):
        paragraphs = []
        layout_metadata = []

        for shape_index, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue

            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                text = paragraph.text.strip()
                if not text:
                    continue

                font_size = _get_font_size(paragraph)
                element_type = _classify_element_type(text, font_size, shape_index, slide_index)

                paragraphs.append(
                    {
                        "shape_id": shape.shape_id,
                        "shape_index": shape_index,
                        "paragraph_index": paragraph_index,
                        "text": text,
                        "font_size": font_size,
                        "position": {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        },
                        "element_type": element_type,
                    }
                )

                layout_metadata.append(
                    {
                        "text_preview": text[:50],
                        "font_size": font_size,
                        "element_type": element_type,
                        "position": "top" if shape.top < 1000000 else "bottom",
                    }
                )

        if paragraphs:
            extracted.append(
                {
                    "slide_id": slide.slide_id,
                    "slide_index": slide_index,
                    "paragraphs": paragraphs,
                    "layout_metadata": layout_metadata,
                    "slide_type": _classify_slide_type(paragraphs),
                }
            )

    return extracted


def _get_font_size(paragraph) -> int:
    """Extract font size in points from a paragraph."""
    for run in paragraph.runs:
        if run.font.size:
            return int(run.font.size.pt)
    return 12


def _classify_element_type(text: str, font_size: int, shape_index: int, slide_index: int) -> str:
    """Classify element type based on font size, text length, and position."""
    word_count = len(text.split())

    if slide_index == 0 and shape_index == 0:
        return "title"
    if font_size >= 40:
        return "title"
    elif font_size >= 28:
        return "heading"
    elif font_size >= 18:
        return "subheading"
    elif word_count <= 4 and font_size >= 20:
        return "heading"
    elif word_count <= 8:
        return "bullet"

    return "body"


def _classify_slide_type(paragraphs: list[dict]) -> str:
    """Classify slide type (title, agenda, body, closing, divider)."""
    if not paragraphs:
        return "blank"

    first_element_type = paragraphs[0].get("element_type", "body")
    paragraph_count = len(paragraphs)

    if first_element_type == "title":
        return "title" if paragraph_count <= 2 else "body"

    first_text = paragraphs[0].get("text", "").lower()
    if any(word in first_text for word in ["agenda", "outline", "contents", "topics"]):
        return "agenda"
    if any(word in first_text for word in ["conclusion", "thank", "questions", "contact"]):
        return "closing"

    return "body"

"""Text replacement in PowerPoint presentations with mode support."""

from pathlib import Path

from pptx import Presentation


def replace_text(
    source_path: Path,
    output_path: Path,
    ai_output: dict[str, list[str]],
    mode: str = "safe"
) -> None:
    """
    Replace text in a presentation with AI-enhanced content.

    Args:
        source_path: Path to input .pptx file
        output_path: Path to save output .pptx file
        ai_output: Dict mapping slide keys to lists of replacement text
        mode: "safe" (1-to-1 mapping) or "creative" (allow restructuring)

    In safe mode: maintains 1-to-1 paragraph mapping.
    In creative mode: intelligently handles bullet count changes.
    """
    presentation = Presentation(str(source_path))

    for slide_number, slide in enumerate(presentation.slides, start=1):
        replacements = ai_output.get(f"slide_{slide_number}")
        if not replacements:
            continue

        if mode == "safe":
            _replace_text_safe(slide, replacements, slide_number)
        else:
            _replace_text_creative(slide, replacements, slide_number)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))


def _replace_text_safe(slide, replacements: list[str], slide_number: int) -> None:
    """
    Replace text maintaining 1-to-1 paragraph mapping (safe mode).
    Skips if replacement count doesn't match original count.
    """
    replacement_index = 0
    original_count = 0

    # First pass: count original paragraphs
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                original_count += 1

    # Safety check
    if len(replacements) != original_count:
        return

    # Second pass: replace
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue

        for paragraph in shape.text_frame.paragraphs:
            if replacement_index >= len(replacements):
                break

            if not paragraph.text.strip():
                continue

            new_text = replacements[replacement_index]
            _set_paragraph_text(paragraph, new_text)
            replacement_index += 1


def _replace_text_creative(slide, replacements: list[str], slide_number: int) -> None:
    """
    Replace text with intelligent restructuring (creative mode).
    Handles cases where bullet count changes.
    """
    if not replacements:
        return

    # Collect all text shapes and their paragraphs
    text_shapes = []
    for shape_index, shape in enumerate(slide.shapes):
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue

        paragraphs_info = []
        for para_index, paragraph in enumerate(shape.text_frame.paragraphs):
            if paragraph.text.strip():
                paragraphs_info.append((para_index, paragraph))

        if paragraphs_info:
            text_shapes.append((shape_index, shape, paragraphs_info))

    if not text_shapes:
        return

    # Distribute replacements across shapes
    # Prefer to replace in the first text shape (main content area)
    main_shape_idx, main_shape, main_paragraphs = text_shapes[0]

    # Clear main shape and fill with all replacements
    _clear_shape_paragraphs(main_shape)
    _fill_shape_with_text(main_shape, replacements)

    # Clear any additional shapes that had original content
    for shape_idx, shape, paragraphs in text_shapes[1:]:
        _clear_shape_paragraphs(shape)


def _clear_shape_paragraphs(shape) -> None:
    """Clear all paragraphs in a shape's text frame."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return

    # Remove all but the first paragraph
    while len(shape.text_frame.paragraphs) > 1:
        # python-pptx doesn't have direct paragraph deletion
        # so we clear text instead
        pass

    # Clear all text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.text = ""


def _fill_shape_with_text(shape, texts: list[str]) -> None:
    """Fill a shape's text frame with a list of texts."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return

    text_frame = shape.text_frame
    # Clear existing text
    for paragraph in text_frame.paragraphs:
        paragraph.text = ""

    for i, text in enumerate(texts):
        if i == 0:
            _set_paragraph_text(text_frame.paragraphs[0], text)
        else:
            new_paragraph = text_frame.add_paragraph()
            _set_paragraph_text(new_paragraph, text)


def _set_paragraph_text(paragraph, new_text: str) -> None:
    """Set paragraph text while preserving runs and formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = new_text
